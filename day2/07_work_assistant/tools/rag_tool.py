"""
RAG 문서 검색 도구
==================
day2/04_RAG/basic_agent/ 코드를 이 프로젝트용으로 재구성했습니다.

핵심 변경 사항:
    - FAISS 인덱스를 미리 만들어두는 대신, 앱 시작 시 docs/ 폴더의 문서를 읽어서 즉시 구축
    - build_vectorstore(docs_dir) : docs_dir의 모든 문서를 파싱 → 벡터스토어 반환
    - create_search_tool(vectorstore) : 벡터스토어를 받아서 LangChain 도구를 반환

사용 방법 (step3, step4에서 이렇게 씁니다):
    from tools.rag_tool import build_vectorstore, create_search_tool

    @st.cache_resource  # 앱 실행 중 한 번만 빌드됩니다
    def load_rag():
        vs = build_vectorstore("docs/")
        return create_search_tool(vs)
"""
import os
from pathlib import Path

import pandas as pd
from pptx import Presentation
from langchain.tools import tool
from langchain_core.documents import Document
from langchain_openai import OpenAIEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter


# ── 문서 파서 함수들 ────────────────────────────────────────────────────────────
# day2/04_RAG/basic_agent/RAG/rag.py의 파서 함수들을 그대로 재활용합니다.

def _parse_excel(file_path: str) -> list[Document]:
    """엑셀 파일의 각 행을 Document로 변환합니다."""
    with pd.ExcelFile(file_path) as xl:
        docs = []
        for sheet_name in xl.sheet_names:
            df = xl.parse(sheet_name, header=1)
            for idx, row in df.iterrows():
                text = "\n".join(f"{col}: {row[col]}" for col in df.columns)
                docs.append(Document(
                    page_content=text,
                    metadata={"source": file_path, "sheet": sheet_name, "row": idx},
                ))
    return docs


def _parse_pdf(file_path: str) -> list[Document]:
    """PDF 파일을 페이지 단위 Document로 변환합니다."""
    return PyPDFLoader(file_path).load()


def _parse_word(file_path: str) -> list[Document]:
    """Word 파일을 Document로 변환합니다."""
    return Docx2txtLoader(file_path).load()


def _parse_pptx(file_path: str) -> list[Document]:
    """PowerPoint 파일을 슬라이드 단위 Document로 변환합니다."""
    prs = Presentation(file_path)
    docs = []
    for slide_num, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        texts.append(para.text.strip())
        if texts:
            docs.append(Document(
                page_content="\n".join(texts),
                metadata={"source": file_path, "slide": slide_num},
            ))
    return docs


def _parse_file(file_path: str) -> list[Document]:
    """확장자에 맞는 파서를 자동으로 선택해서 Document 리스트를 반환합니다."""
    ext = Path(file_path).suffix.lower()
    if ext in (".doc", ".docx"):
        return _parse_word(file_path)
    elif ext == ".xlsx":
        return _parse_excel(file_path)
    elif ext == ".pdf":
        return _parse_pdf(file_path)
    elif ext in (".ppt", ".pptx"):
        return _parse_pptx(file_path)
    return []  # 지원하지 않는 형식은 건너뜁니다


# ── 공개 함수 ─────────────────────────────────────────────────────────────────

def build_vectorstore(docs_dir: str, chunk_size: int = 500, chunk_overlap: int = 50) -> FAISS:
    """
    docs_dir 안의 모든 문서를 파싱하고 FAISS 벡터스토어를 구축합니다.

    Args:
        docs_dir: 문서들이 있는 폴더 경로 (예: "docs/")
        chunk_size: 텍스트를 자르는 단위 (기본 500자)
        chunk_overlap: 청크 간 겹치는 문자 수 (문맥 유지를 위해)

    Returns:
        검색 가능한 FAISS 벡터스토어

    Notes:
        이 함수는 OpenAI Embeddings API를 호출하므로 시간이 걸립니다.
        Streamlit에서는 반드시 @st.cache_resource로 감싸서 한 번만 실행하세요.
    """
    # 1. docs_dir 안의 모든 파일을 파싱해서 Document 리스트 생성
    all_docs = []
    for file_name in os.listdir(docs_dir):
        file_path = os.path.join(docs_dir, file_name)
        docs = _parse_file(file_path)
        all_docs.extend(docs)
        print(f"  [{file_name}] → {len(docs)}개 청크")

    if not all_docs:
        raise ValueError(f"'{docs_dir}' 폴더에서 파싱 가능한 문서를 찾을 수 없습니다.")

    # 2. 긴 문서를 적절한 크기의 청크(덩어리)로 자릅니다
    # 청크가 너무 크면 검색 정확도가 떨어지고, 너무 작으면 문맥이 끊깁니다
    splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    chunks = splitter.split_documents(all_docs)
    print(f"총 {len(chunks)}개 청크로 분할 완료")

    # 3. 각 청크를 벡터(숫자 배열)로 변환하고 FAISS에 저장합니다
    # text-embedding-3-small: 빠르고 저렴한 OpenAI 임베딩 모델
    vectorstore = FAISS.from_documents(chunks, OpenAIEmbeddings(model="text-embedding-3-small"))
    print("벡터스토어 구축 완료!")

    return vectorstore


def create_search_tool(vectorstore: FAISS):
    """
    벡터스토어를 받아서 LangChain 도구를 생성해 반환합니다.

    왜 함수로 감쌌나?
        @tool 데코레이터는 함수 정의 시점에 고정됩니다.
        vectorstore를 외부에서 주입받으려면 이렇게 클로저(closure) 패턴을 씁니다.
        덕분에 vectorstore를 전역 변수로 쓰지 않아도 됩니다.
    """

    @tool
    def search_documents(query: str) -> str:
        """
        키키테크 사내 문서에서 정보를 검색합니다.
        회사 소개, 제품 카탈로그, 사내 규정, 임직원 및 프로젝트 현황을 확인할 수 있습니다.
        """
        # 의미적으로 유사한 상위 3개 청크를 검색합니다
        results = vectorstore.similarity_search(query, k=3)

        formatted = []
        for doc in results:
            # 출처 파일명만 추출 (전체 경로 대신 파일명만 보여줘서 가독성 향상)
            source = Path(doc.metadata.get("source", "unknown")).name
            formatted.append(f"[출처: {source}]\n{doc.page_content}")

        return "\n\n---\n\n".join(formatted)

    return search_documents
