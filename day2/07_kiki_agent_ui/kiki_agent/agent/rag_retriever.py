import os
import requests
import pandas as pd
from pathlib import Path
from pptx import Presentation
from dotenv import load_dotenv

from langchain_core.documents import Document
from langchain_core.embeddings import Embeddings
from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader
from langchain_community.retrievers import BM25Retriever
from langchain_community.vectorstores import FAISS
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_classic.retrievers import EnsembleRetriever


load_dotenv(dotenv_path="../../.env")


class MyCustomEmbeddings(Embeddings):
    def __init__(self, api_url, ticket):
        self.api_url = api_url
        self.ticket = ticket

    def embed_documents(self, texts):
        response = requests.post(
            f"{self.api_url}/v1/embeddings",
            json={"input": texts},
            headers={"Content-Type": "application/json", "x-dep-ticket": self.ticket}
        )
        return [item["embedding"] for item in response.json()['data']]

    def embed_query(self, text):
        return self.embed_documents([text])[0]


class RAGRetriever:
    def __init__(self, chunk_size=500, chunk_overlap=50):
        # 임베딩 모델 객체 생성
        self.custom_embedding_model = MyCustomEmbeddings(
            api_url=os.getenv("embedding_api_url"),
            ticket=os.getenv("embedding_x_dep_ticket"),
        )

        # chunk_size = 한 청크의 최대 글자 수 / chunk_overlap = 겹치는 영역의 글자 수
        self.splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)

        self.total_chunks = []
        self.ensemble_retriever = None

    # ── 파싱 함수들 ───────────────────────────────────────────────────────────

    def parse_excel(self, file_path):
        with pd.ExcelFile(file_path) as xl:
            docs = []
            for sheet_name in xl.sheet_names:
                df = xl.parse(sheet_name, header=1)
                for idx, row in df.iterrows():
                    text = "\n".join(f"{col}: {row[col]}" for col in df.columns)
                    docs.append(Document(
                        page_content=text,
                        metadata={
                            "source": file_path,
                            "sheet": sheet_name,
                            "row": idx,
                        }
                    ))
        return docs

    def parse_pdf(self, file_path):
        docs = PyPDFLoader(file_path).load()
        return docs

    def parse_word(self, file_path):
        docs = Docx2txtLoader(file_path).load()
        return docs

    def parse_pptx(self, file_path):
        prs = Presentation(file_path)
        docs = []
        for slide_num, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            texts.append(para.text.strip())
            if texts:
                docs.append(Document(
                    page_content="\n".join(texts),
                    metadata={
                        "source": file_path,
                        "slide": slide_num,
                    }
                ))
        return docs

    # ── 파일 추가 ─────────────────────────────────────────────────────────────

    def add_file(self, file_path):
        ext = Path(file_path).suffix.lower()

        if ext in (".xlsx", ".xls"):
            docs = self.parse_excel(file_path)
        elif ext == ".pdf":
            docs = self.parse_pdf(file_path)
        elif ext in (".docx", ".doc"):
            docs = self.parse_word(file_path)
        elif ext in (".pptx", ".ppt"):
            docs = self.parse_pptx(file_path)

        chunks = self.splitter.split_documents(docs)
        self.total_chunks += chunks

    # ── EnsembleRetriever 생성 ────────────────────────────────────────────────

    def build_retriever(self, sparse_weight=0.4, dense_weight=0.6, k=5):
        # BM25 sparse retriever 생성
        bm25_retriever = BM25Retriever.from_documents(self.total_chunks)
        bm25_retriever.k = k

        # 데이터 준비
        page_contents = [chunk.page_content for chunk in self.total_chunks]
        page_metadatas = [chunk.metadata for chunk in self.total_chunks]

        # 임베딩 생성
        embeddings = self.custom_embedding_model.embed_documents(page_contents)

        # FAISS 벡터스토어 생성
        text_embeddings = list(zip(page_contents, embeddings))
        vectorstore = FAISS.from_embeddings(text_embeddings, self.custom_embedding_model, metadatas=page_metadatas)

        # FAISS dense retriever
        dense_retriever = vectorstore.as_retriever(search_kwargs={"k": k})

        # Sparse + Dense 비율로 혼합
        self.ensemble_retriever = EnsembleRetriever(
            retrievers=[bm25_retriever, dense_retriever],
            weights=[sparse_weight, dense_weight]
        )

    def search(self, query):
        results = self.ensemble_retriever.invoke(query)
        return results
