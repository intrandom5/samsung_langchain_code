import pandas as pd
from pathlib import Path
from pptx import Presentation
from langchain_core.documents import Document
from langchain_community.document_loaders import PyPDFLoader
from langchain_community.document_loaders import Docx2txtLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter


def parse_excel(file_path: str) -> list[Document]:
    with pd.ExcelFile(file_path) as xl:
        docs = []      
        for sheet_name in xl.sheet_names:
            df = xl.parse(sheet_name, header=1)
            for idx, row in df.iterrows():
                text = "\n".join(f"{col}: {row[col]}" for col in df.columns)
                docs.append(Document(
                    page_content=text, 
                    metadata={
                        "source": file_path, 
                        "sheet": sheet_name, 
                        "row": idx, 
                    }
                ))
    return docs

def parse_pdf(file_path: str) -> list[Document]:
    docs = PyPDFLoader(file_path).load()
    return docs

def parse_word(file_path: str) -> list[Document]:
    # 전달받은 경로(file_path)를 사용하도록 수정했습니다.
    docs = Docx2txtLoader(file_path).load()
    return docs

def parse_pptx(file_path: str) -> list[Document]:
    prs = Presentation(file_path)
    docs = []
    for slide_num, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        texts.append(para.text.strip())
        if texts:
            docs.append(Document(
                page_content="\n".join(texts),
                metadata={
                    "source": file_path,
                    "slide": slide_num,
                }
            ))
    return docs

# chunk_size = 한 청크의 최대 글자 수 / chunk_overlap = 겹치는 영역의 글자 수
def chunk_splitter(docs: list[Document], chunk_size: int = 500, chunk_overlap: int = 50) -> list[Document]:
    # 파이썬 문법에 맞게 인자 선언 방식을 수정하고, 전달받은 값을 사용하도록 고쳤습니다.
    splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    chunks = splitter.split_documents(docs)
    return chunks