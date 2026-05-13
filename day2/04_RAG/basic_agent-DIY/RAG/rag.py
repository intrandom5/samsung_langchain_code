import os, shutil
import pandas as pd
from pathlib import Path
from pptx import Presentation
from dotenv import load_dotenv

from langchain_openai import OpenAIEmbeddings
from langchain_core.documents import Document
from langchain_community.vectorstores import FAISS
from langchain_community.document_loaders import PyPDFLoader
from langchain_community.document_loaders import Docx2txtLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter


load_dotenv(dotenv_path="../.env")

embeddings = OpenAIEmbeddings(
    model="text-embedding-3-small",
    api_key=os.getenv("OPENAI_API_KEY")
)


def parse_excel(file_path: str) -> list[Document]:
    # 1. 전달받은 경로(file_path)에 있는 엑셀 파일을 읽어옵니다.
    with pd.ExcelFile(file_path) as xl:
    # 2. 데이터를 담아둘 빈 바구니(리스트)를 만듭니다.
        docs = []
        
        # 3. 엑셀 파일 안에 있는 여러 개의 '시트(sheet)'를 하나씩 꺼내서 반복합니다.
        for sheet_name in xl.sheet_names:
            # 4. 현재 시트의 내용을 표(데이터프레임) 형태로 불러옵니다.
            # header=1 은 엑셀의 두 번째 줄(행)을 제목(컬럼 이름)으로 사용한다는 뜻입니다.
            df = xl.parse(sheet_name, header=1)
            
            # 5. 표의 데이터를 한 줄씩 꺼내서 반복합니다.
            for idx, row in df.iterrows():
                # 6. 한 줄에 있는 각 칸의 데이터를 "제목: 내용" 형식으로 묶어서 하나의 긴 글로 만듭니다.
                # (예: "이름: 홍길동 \n 나이: 20")
                text = "\n".join(f"{col}: {row[col]}" for col in df.columns)
                
                # 7. 만들어진 글과 출처 정보를 하나로 묶어서(Document 객체) 바구니에 담습니다.
                docs.append(Document(
                    page_content=text, # 실제 내용이 들어가는 부분
                    metadata={         # 이 데이터가 어디서 왔는지 기록하는 부분
                        "source": file_path, # 어떤 파일에서 왔는지
                        "sheet": sheet_name, # 어떤 시트에서 왔는지
                        "row": idx,          # 몇 번째 줄(행)에서 왔는지
                    }
                ))
                
    # 8. 모든 데이터가 담긴 바구니(문서 리스트)를 최종적으로 반환합니다.
    return docs


def parse_pdf(file_path: str) -> list[Document]:
    # pdf 파일을 불러오는 코드
    pdf_contents = PyPDFLoader(file_path).load()

    return pdf_contents

# word를 입력 받아 텍스트로 변환하는 함수
def parse_word(file_path: str) -> list[Document]:
    # word 파일을 불러오는 코드
    word_contents = Docx2txtLoader(file_path).load()

    return word_contents


# ppt를 입력 받아 텍스트로 변환하는 함수
def parse_pptx(file_path: str) -> list[Document]:
    # ppt 파일 로드
    prs = Presentation(file_path)
    docs = []
    # 각 슬라이드를 돌며...
    for slide_num, slide in enumerate(prs.slides, 1):
        texts = []
        # 슬라이드의 모든 요소를 살펴보며 텍스트가 있으면 texts 리스트에 추가.
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        texts.append(para.text.strip())
        if texts:
            docs.append(Document(
                page_content="\n".join(texts),
                metadata={
                    "source": file_path,
                    "slide": slide_num,
                }
            ))
    return docs

def split_documents(docs: list[Document], chunk_size: int, chunk_overlap: int) -> list[Document]:
    splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    chunks = splitter.split_documents(docs)
    return chunks

def register_rag_docs(new_docs: str = "new_docs", chunk_size: int = 500, chunk_overlap: int = 50):
    new_files = os.listdir(new_docs)
    print(f"총 {len(new_files)}개의 파일을 찾았습니다!")

    total_docs = []

    for new_file in new_files:
        file_path = os.path.join("new_docs", new_file)
        if file_path.endswith(".doc") or file_path.endswith(".docs") or file_path.endswith(".docx"):
            docs = parse_word(file_path)
        elif file_path.endswith(".xlsx"):
            docs = parse_excel(file_path)
        elif file_path.endswith(".pdf"):
            docs = parse_pdf(file_path)
        elif file_path.endswith(".ppt") or file_path.endswith(".pptx"):
            docs = parse_pptx(file_path)
        else:
            continue
        total_docs += docs

        registered_path = os.path.join("registered_docs", new_file)
        shutil.move(file_path, registered_path)
    
    chunks = split_documents(total_docs, chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    print(f"총 {len(chunks)}개의 문서가 생성되었습니다!")
    vectorstore = FAISS.from_documents(chunks, OpenAIEmbeddings())

    # 이미 기존의 faiss 벡터가 존재한다면, 기존의 벡터에 추가
    if os.path.exists("faiss_index"):
        original_vectorstore = FAISS.load_local(
            "faiss_index",
            OpenAIEmbeddings(),
            allow_dangerous_deserialization=True
        )
        vectorstore.merge_from(original_vectorstore)

    vectorstore.save_local("faiss_index")

def main():
    register_rag_docs(chunk_size=500, chunk_overlap=50)

if __name__ == "__main__":
    main()